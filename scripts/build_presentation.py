"""Generate the project presentation as a PowerPoint file.

Run::

    .venv/bin/python scripts/build_presentation.py

Output lands at ``presentation.pptx`` in the project root. Re-runnable;
overwrites the previous file. Edit this script (not the .pptx) when
slides need structural changes.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentation.pptx"
SCREENSHOT_DIR = ROOT / "docs" / "screenshots"

# --- design tokens ---------------------------------------------------
INK = RGBColor(0x1F, 0x21, 0x28)
INK_SOFT = RGBColor(0x4F, 0x54, 0x5C)
INK_MUTED = RGBColor(0x8B, 0x8D, 0x98)
ACCENT = RGBColor(0x5E, 0x6A, 0xD2)  # the Linear-ish purple used in the UI
ACCENT_SOFT = RGBColor(0xEC, 0xEE, 0xFB)
GREEN = RGBColor(0x30, 0xA4, 0x6C)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
RED = RGBColor(0xE5, 0x48, 0x4D)
BG_PAGE = RGBColor(0xFA, 0xFB, 0xFC)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE4, 0xE8)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _set_slide_background(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align=None,
    font: str = "Helvetica Neue",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def _add_bullets(
    slide,
    items: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 16,
    color: RGBColor = INK,
    bullet_color: RGBColor = ACCENT,
    line_spacing: float = 1.35,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        # bullet
        b = p.add_run()
        b.text = "▸  "
        b.font.name = "Helvetica Neue"
        b.font.size = Pt(size)
        b.font.color.rgb = bullet_color
        b.font.bold = True
        # body
        r = p.add_run()
        r.text = item
        r.font.name = "Helvetica Neue"
        r.font.size = Pt(size)
        r.font.color.rgb = color


def _add_rect(slide, *, left, top, width, height, fill=BG_WHITE, line=BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.adjustments[0] = 0.08  # corner radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _set_shape_text(shape, text: str, *, size: int = 14, color: RGBColor = INK, bold: bool = False, align=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.14)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    p.text = text
    run = p.runs[0]
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _add_footer(slide, page_num: int, total: int, title: str = "AI-Assisted Ticket Triage"):
    _add_text(slide, title, left=0.5, top=7.05, width=7, height=0.4,
              size=10, color=INK_MUTED)
    _add_text(slide, f"{page_num} / {total}", left=12.0, top=7.05, width=1.0, height=0.4,
              size=10, color=INK_MUTED)


def _header(slide, eyebrow: str, title: str):
    _add_text(slide, eyebrow.upper(), left=0.6, top=0.5, width=12, height=0.4,
              size=12, color=ACCENT, bold=True)
    _add_text(slide, title, left=0.6, top=0.9, width=12, height=0.9,
              size=32, color=INK, bold=True)
    # accent line under header
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.7), Inches(1.2), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


# ---------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------

def slide_title(prs, total_pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)

    # Eyebrow
    _add_text(slide, "FALLSTUDIE", left=1.0, top=2.0, width=10, height=0.4,
              size=14, color=ACCENT, bold=True)

    # Title
    _add_text(slide, "AI-Assisted Ticket Triage Platform",
              left=1.0, top=2.5, width=11, height=1.4,
              size=48, color=INK, bold=True)

    # Subtitle
    _add_text(
        slide,
        "Retrieval-Augmented Triage + Knowledge-Engineering Prioritization\n"
        "für den Bank-IT-Support",
        left=1.0, top=4.0, width=11, height=1.2,
        size=22, color=INK_SOFT,
    )

    # Meta line
    _add_text(slide, "FastAPI · React · scikit-learn · LiteLLM · Render",
              left=1.0, top=5.4, width=11, height=0.4,
              size=14, color=INK_MUTED)
    _add_text(slide, "Claudio Vinci  ·  2026",
              left=1.0, top=5.85, width=11, height=0.4,
              size=14, color=INK_MUTED)


def slide_use_case(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Problem & Use Case", "Bank-IT-Support ertrinkt in unstrukturierten Tickets")

    _add_bullets(slide, [
        "Operator:innen erhalten täglich 50–200 Tickets — VPN, Passwort, SEPA, Mobile-App, Compliance.",
        "Klassische „AI-Triage“ versteckt ihre Begründung oder ignoriert historische Routing-Entscheidungen.",
        "Lösung: Empfehlung mit Begründung + Top-3 ähnliche bereits gelöste Fälle + KE-Priorisierung.",
        "Operator bleibt in Kontrolle: AI schlägt vor, Mensch akzeptiert/überschreibt, Entscheidung trainiert das System weiter.",
    ], left=0.6, top=2.0, width=8.5, height=4.0, size=17)

    # Right-side stat card
    _add_rect(slide, left=9.4, top=2.0, width=3.4, height=2.3, fill=ACCENT_SOFT, line=ACCENT)
    _add_text(slide, "16 %", left=9.6, top=2.15, width=3.2, height=0.9,
              size=44, color=ACCENT, bold=True)
    _add_text(slide, "der Tickets im Demo-Korpus sind",
              left=9.6, top=3.05, width=3.2, height=0.4, size=13, color=INK_SOFT)
    _add_text(slide, "Self-Service-fähig", left=9.6, top=3.45, width=3.2, height=0.4,
              size=15, color=INK, bold=True)
    _add_text(slide, "→ Runbook statt Queue", left=9.6, top=3.85, width=3.2, height=0.3,
              size=12, color=INK_MUTED)

    _add_rect(slide, left=9.4, top=4.55, width=3.4, height=2.0, fill=ACCENT_SOFT, line=ACCENT)
    _add_text(slide, "99", left=9.6, top=4.7, width=3.2, height=0.9,
              size=44, color=ACCENT, bold=True)
    _add_text(slide, "Tickets im Live-Korpus,",
              left=9.6, top=5.6, width=3.2, height=0.4, size=13, color=INK_SOFT)
    _add_text(slide, "100 % KE-priorisiert", left=9.6, top=6.0, width=3.2, height=0.4,
              size=15, color=INK, bold=True)

    _add_footer(slide, n, total_pages)


def slide_architecture(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Architektur", "Hexagonal — Dependency-Arrows zeigen einwärts")

    # Layer boxes from outside to inside
    layers = [
        ("Interfaces (FastAPI)", "routes · schemas · middleware · error handlers · rate-limit", ACCENT_SOFT),
        ("Application", "use cases · ports · DTOs", RGBColor(0xFD, 0xF6, 0xE3)),
        ("Domain", "entities · enums · rules · pure Python", RGBColor(0xE8, 0xF4, 0xEE)),
        ("Infrastructure (Adapters)", "SQLAlchemy · LiteLLM · scikit-learn · YAML policy · logging", RGBColor(0xEC, 0xF1, 0xF5)),
    ]
    y = 2.0
    for label, sub, color in layers:
        _add_rect(slide, left=0.6, top=y, width=7.5, height=1.0, fill=color, line=BORDER)
        _add_text(slide, label, left=0.85, top=y + 0.1, width=7, height=0.4,
                  size=17, color=INK, bold=True)
        _add_text(slide, sub, left=0.85, top=y + 0.5, width=7, height=0.4,
                  size=12, color=INK_SOFT)
        y += 1.2

    # Right side: contract bullets
    _add_text(slide, "Vertrag", left=8.6, top=2.0, width=4.5, height=0.5,
              size=18, color=INK, bold=True)
    _add_bullets(slide, [
        "Domain importiert kein Framework.",
        "Application kennt nur Ports, keine Adapter.",
        "4 Ports: Classifier · SimilarTickets · Prioritization · TicketRepository.",
        "Swap von SQLite→Postgres oder LiteLLM→anderer LLM: 1 Env-Var bzw. 1 Datei.",
        "12 Use Cases · 7 Entitäten · alles per Pydantic-Schema an die API gebunden.",
        "4 ADRs dokumentieren architektonische Entscheidungen (docs/adr/).",
    ], left=8.6, top=2.55, width=4.5, height=4.5, size=13)

    _add_footer(slide, n, total_pages)


def slide_domain(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Domänenmodell", "7 Entitäten · 12 Use Cases · 4 Ports")

    # Entities grid
    entities = [
        ("Ticket", "Kern-Aggregate"),
        ("TriageAnalysis", "Output des Klassifikators"),
        ("Prioritization", "KE-Bewertung (neu)"),
        ("TriageDecision", "Review-Entscheidung"),
        ("Assignment", "Team/Bearbeitung"),
        ("SimilarCase", "RAG-Nachbar (read-only)"),
        ("TicketEvent", "Audit-Trail"),
    ]
    cols = 2
    box_w, box_h = 5.6, 0.75
    for i, (name, sub) in enumerate(entities):
        r, c = divmod(i, cols)
        x = 0.6 + c * (box_w + 0.4)
        y = 2.0 + r * (box_h + 0.18)
        _add_rect(slide, left=x, top=y, width=box_w, height=box_h, fill=BG_WHITE, line=BORDER)
        _add_text(slide, name, left=x + 0.18, top=y + 0.08, width=box_w - 0.36, height=0.35,
                  size=14, color=ACCENT, bold=True)
        _add_text(slide, sub, left=x + 0.18, top=y + 0.4, width=box_w - 0.36, height=0.3,
                  size=11, color=INK_SOFT)

    _add_text(slide, "Use Cases (Auszug)", left=0.6, top=5.65, width=8, height=0.4,
              size=15, color=INK, bold=True)
    _add_text(slide,
              "triage_ticket  ·  save_triage_decision  ·  assign_ticket  ·  escalate_ticket  ·  "
              "update_ticket_status  ·  backfill_prioritization  ·  get_dashboard_analytics",
              left=0.6, top=6.05, width=12, height=0.8, size=12, color=INK_SOFT)

    _add_footer(slide, n, total_pages)


def slide_ai_layers(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "AI-Pipeline", "Drei Schichten — Klassifikation × Retrieval × Wissensregeln")

    cards = [
        ("1.  LLM-Klassifikation",
         "LiteLLM-Proxy mit Azure GPT-OSS-120B.\n"
         "Strukturierte JSON-Ausgabe: category, priority, suggested_team, summary, next_step, rationale.\n"
         "Recovery für abgeschnittene Completions.",
         ACCENT),
        ("2.  Retrieval (RAG)",
         "scikit-learn TF-IDF + Cosine NearestNeighbors.\n"
         "Korpus: ausschließlich Tickets mit reviewed_by IS NOT NULL.\n"
         "Top-3 ähnliche Fälle werden dem LLM als Kontext UND dem Operator in der UI angezeigt.",
         GREEN),
        ("3.  Knowledge-Engineering",
         "YAML-Policy mit 17 Regeln in 4 Tiers.\n"
         "Bewertet jedes Ticket auf Wichtigkeit · Dringlichkeit · Aufwand · Lösbarkeit.\n"
         "Aufwand wird aus RAG-Nachbarn gemittelt (lernend), Regeln sind explainable.",
         AMBER),
    ]
    y = 2.0
    for title, body, accent in cards:
        _add_rect(slide, left=0.6, top=y, width=12.2, height=1.55, fill=BG_WHITE, line=accent)
        _add_text(slide, title, left=0.85, top=y + 0.1, width=11.5, height=0.4,
                  size=18, color=accent, bold=True)
        _add_text(slide, body, left=0.85, top=y + 0.5, width=11.7, height=1.0,
                  size=13, color=INK_SOFT)
        y += 1.7

    _add_footer(slide, n, total_pages)


def slide_ke_detail(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Knowledge-Engineering Layer", "Prio-Score = Wichtigkeit × Dringlichkeit  (1–25)")

    # Two columns: dimensions table + sample YAML
    headers = [("Dimension", "Antwort auf", "Quelle")]
    rows = [
        ("Wichtigkeit (1–5)", "Wie kritisch ist das?", "YAML-Regel"),
        ("Dringlichkeit (1–5)", "Wie schnell?", "YAML-Regel"),
        ("Aufwand (Min.)", "Wie teuer?", "Ø RAG-Nachbarn"),
        ("Lösbarkeit", "Wer muss ran?", "YAML-Regel"),
        ("Prio-Score (1–25)", "Backlog-Sort", "Wichtigkeit × Dringlichkeit"),
        ("Auto-Resolve?", "Self-Service möglich?", "self-service ∧ confidence > 0.6"),
    ]
    box_x, box_y, box_w, box_h = 0.6, 2.0, 6.5, 4.4
    _add_rect(slide, left=box_x, top=box_y, width=box_w, height=box_h, fill=BG_WHITE, line=BORDER)
    rh = box_h / (len(rows) + 1)
    for i, (a, b, c) in enumerate(headers + rows):
        ry = box_y + i * rh
        if i == 0:
            _add_text(slide, a, left=box_x + 0.1, top=ry + 0.05, width=2.4, height=rh,
                      size=12, color=INK_MUTED, bold=True)
            _add_text(slide, b, left=box_x + 2.4, top=ry + 0.05, width=2.0, height=rh,
                      size=12, color=INK_MUTED, bold=True)
            _add_text(slide, c, left=box_x + 4.4, top=ry + 0.05, width=2.0, height=rh,
                      size=12, color=INK_MUTED, bold=True)
        else:
            _add_text(slide, a, left=box_x + 0.1, top=ry + 0.05, width=2.4, height=rh,
                      size=12, color=INK, bold=True)
            _add_text(slide, b, left=box_x + 2.4, top=ry + 0.05, width=2.0, height=rh,
                      size=11, color=INK_SOFT)
            _add_text(slide, c, left=box_x + 4.4, top=ry + 0.05, width=2.0, height=rh,
                      size=11, color=INK_SOFT)

    # YAML snippet on the right
    yaml_box = _add_rect(slide, left=7.5, top=2.0, width=5.3, height=4.4,
                        fill=RGBColor(0x16, 0x18, 0x21), line=RGBColor(0x16, 0x18, 0x21))
    yaml_text = (
        "rules:\n"
        "  - id: aml-critical\n"
        "    match:\n"
        "      tags_any: [aml, sanctions, kyc]\n"
        "      title_any: [aml, sanktion, kyc]\n"
        "    set:\n"
        "      impact_score: 5\n"
        "      urgency_score: 5\n"
        "      solvability: specialist\n"
        "      rationale: \"AML — sofortige\n"
        "        Eskalation.\"\n"
        "default:\n"
        "  impact_score: 2\n"
        "  urgency_score: 2\n"
        "  solvability: l2"
    )
    tf = yaml_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = yaml_text
    for run in p.runs:
        run.font.name = "Menlo"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0xE6, 0xE6, 0xE6)

    _add_footer(slide, n, total_pages)


def slide_coding_agents(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "AI-Coding-Agents im Einsatz  (Bewertungsgewicht: 40 %)",
            "Reflektierter Spec-first-Workflow mit Claude Code")

    _add_bullets(slide, [
        "Primary Agent: Claude Code (Anthropic), Modell Opus 4.7 — Hexagonal-Architektur-fest, kontextstark.",
        "Spec-first: jede neue Funktion startet mit 1-Absatz-Vertrag (Input/Output, Ports, Entitäten).",
        "TDD-Recipe in AGENTS.md kodifiziert: Test schreiben → rot → Use Case → Adapter → Route → Frontend.",
        "Explore-Subagents für Codebase-Maps vor jedem größeren Eingriff (parallel, schont Kontext).",
        "Plan-Mode für nicht-triviale Änderungen — Approval bevor geschrieben wird.",
        "Hard-Rules in AGENTS.md verhindern Stack-Drift (keine Domain-Framework-Imports, FK-Cascade, …).",
        "User-Korrekturen werden als operative Regeln in AGENTS.md zementiert — Vermeidung von Wiederholungen.",
    ], left=0.6, top=2.0, width=8.5, height=4.6, size=15)

    # Side: stat tiles
    _add_rect(slide, left=9.4, top=2.0, width=3.4, height=1.3, fill=ACCENT_SOFT, line=ACCENT)
    _add_text(slide, "10", left=9.6, top=2.1, width=3.2, height=0.6,
              size=40, color=ACCENT, bold=True)
    _add_text(slide, "Hard-Rules in AGENTS.md", left=9.6, top=2.85, width=3.2, height=0.3,
              size=11, color=INK_SOFT)

    _add_rect(slide, left=9.4, top=3.5, width=3.4, height=1.3, fill=ACCENT_SOFT, line=ACCENT)
    _add_text(slide, "4", left=9.6, top=3.6, width=3.2, height=0.6,
              size=40, color=ACCENT, bold=True)
    _add_text(slide, "ADRs für nicht-triviale Entscheidungen", left=9.6, top=4.35, width=3.2, height=0.4,
              size=11, color=INK_SOFT)

    _add_rect(slide, left=9.4, top=5.0, width=3.4, height=1.3, fill=ACCENT_SOFT, line=ACCENT)
    _add_text(slide, "8", left=9.6, top=5.1, width=3.2, height=0.6,
              size=40, color=ACCENT, bold=True)
    _add_text(slide, "Real-bug Pitfalls dokumentiert", left=9.6, top=5.85, width=3.2, height=0.3,
              size=11, color=INK_SOFT)

    _add_footer(slide, n, total_pages)


def slide_agents_md(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "AGENTS.md", "Spec-first-Artefakt — der Agent-Vertrag in einer Datei")

    _add_text(slide, "Inhalt (29 Sections)",
              left=0.6, top=2.0, width=6, height=0.4, size=16, color=INK, bold=True)
    _add_bullets(slide, [
        "60-Sekunden-Codebase-Summary",
        "10 Hard-Rules (Architektur, RAG-Filter, FK-Cascade, …)",
        "Critical Commands — exakte CLI-Aufrufe",
        "Code-Style: Python (Dataclass/Pydantic), JS (kein TS)",
        "Conventional-Commits-Template mit HEREDOC",
        "Architecture-Contract: Ports-Tabelle",
        "Domain-Knowledge: KE-Layer, Seed-Korpus, Migrations",
        "Workflow für neue Features (12-Schritt-Rezept)",
        "Common Pitfalls — 8 reale Bugs aus diesem Repo",
        "When to ask vs. proceed — Entscheidungsmatrix",
        "„Files to read first“-Lookup je Tasktyp",
    ], left=0.6, top=2.45, width=6, height=4.6, size=12)

    # Right: quote block from AGENTS.md
    _add_rect(slide, left=7.0, top=2.0, width=5.8, height=5.0, fill=BG_WHITE, line=ACCENT)
    _add_text(slide, "Auszug — Hard Rule #6", left=7.2, top=2.1, width=5.4, height=0.4,
              size=12, color=ACCENT, bold=True)
    _add_text(slide,
              "„Postgres enforces foreign keys; SQLite doesn't. "
              "When deleting tickets, always delete ticket_events first. "
              "Use _delete_tickets() — a direct DELETE will succeed on SQLite "
              "and fail on Postgres with ForeignKeyViolation.“",
              left=7.2, top=2.55, width=5.4, height=2.0,
              size=13, color=INK)
    _add_text(slide, "Realer Bug in dieser Codebasis · Fix dokumentiert",
              left=7.2, top=4.7, width=5.4, height=0.3,
              size=11, color=INK_MUTED)

    _add_text(slide, "Pre-empts: User wiederholt mich nicht.",
              left=7.2, top=5.6, width=5.4, height=0.4,
              size=14, color=INK_SOFT, bold=True)
    _add_text(slide,
              "Jede User-Korrektur die zukünftig relevant ist landet in AGENTS.md — der nächste Agent-Lauf liest sie und vermeidet den Fehler.",
              left=7.2, top=6.0, width=5.4, height=1.0,
              size=11, color=INK_SOFT)

    _add_footer(slide, n, total_pages)


def slide_tools(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Toolauswahl & Begründung", "Jede Wahl trade-off-getrieben, nicht Hype")

    headers = ("Bereich", "Tool", "Warum")
    rows = [
        ("API", "FastAPI", "Async-fähig, Pydantic-Integration, OpenAPI gratis"),
        ("ORM", "SQLAlchemy 2.0", "Dual-Backend SQLite/Postgres, mature, gut testbar"),
        ("LLM-Gateway", "LiteLLM", "Provider-agnostisch — Lock-in vermeiden (ADR 0003)"),
        ("Klassifikation", "scikit-learn TF-IDF + NB", "Klein, schnell, deterministisch — kein Embedding-Overhead bei <10k Tickets"),
        ("RAG", "scikit-learn TfidfVectorizer", "Gleiche Begründung wie oben — sentence-transformers wäre Over-engineering (ADR 0004)"),
        ("KE-Engine", "Eigener YAML-Adapter", "Operatoren editieren Regeln, kein Code-Change"),
        ("Frontend", "React + Vite + Recharts", "Schneller Dev-Cycle, kleine Bundles, Charts ohne extra Lib"),
        ("Hosting", "Render (Blueprint)", "Free Tier reicht für Demo, Auto-Deploy via Hook"),
        ("Quality", "ruff · pytest · Vitest · Playwright · bandit", "Eine CI-Pipeline, alle Layers"),
    ]
    box_x, box_y, box_w = 0.6, 2.0, 12.2
    rh = 0.43
    # Header row
    _add_rect(slide, left=box_x, top=box_y, width=box_w, height=rh,
              fill=ACCENT_SOFT, line=BORDER)
    cols_x = [box_x + 0.15, box_x + 1.7, box_x + 4.6]
    for i, h in enumerate(headers):
        _add_text(slide, h, left=cols_x[i], top=box_y + 0.05, width=4, height=rh,
                  size=12, color=INK, bold=True)
    # Body rows
    for i, row in enumerate(rows):
        ry = box_y + (i + 1) * rh
        bg = BG_WHITE if i % 2 == 0 else BG_PAGE
        _add_rect(slide, left=box_x, top=ry, width=box_w, height=rh,
                  fill=bg, line=BORDER)
        for j, cell in enumerate(row):
            _add_text(slide, cell, left=cols_x[j], top=ry + 0.05, width=4.5, height=rh,
                      size=11, color=INK if j < 2 else INK_SOFT, bold=(j == 1))

    _add_footer(slide, n, total_pages)


def slide_testing(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Test-Strategie", "Drei-Schicht-Pyramide · 75 %-Coverage-Gate · CI-blocking")

    # Pyramid visualization on the left
    levels = [
        ("E2E (Playwright)", "1–2 happy-path flows", RED),
        ("API (FastAPI TestClient)", "HTTP-Schemas, Auth, Rate-Limits, Mocked LLM", AMBER),
        ("Application (Use Cases)", "Coordinator-Logik gegen Fakes", ACCENT),
        ("Unit (Pure Python)", "Domain-Regeln, Policy-Matcher, Cache, …", GREEN),
    ]
    base_x = 0.8
    widths = [3.4, 5.0, 6.0, 7.0]
    y = 2.1
    for (label, sub, color), w in zip(levels, widths):
        x = base_x + (7.0 - w) / 2
        _add_rect(slide, left=x, top=y, width=w, height=0.85, fill=BG_WHITE, line=color)
        _add_text(slide, label, left=x + 0.15, top=y + 0.08, width=w - 0.3, height=0.35,
                  size=13, color=color, bold=True)
        _add_text(slide, sub, left=x + 0.15, top=y + 0.42, width=w - 0.3, height=0.35,
                  size=11, color=INK_SOFT)
        y += 1.0

    # Right side: rules & numbers
    _add_text(slide, "Regeln", left=8.6, top=2.1, width=4, height=0.4,
              size=16, color=INK, bold=True)
    _add_bullets(slide, [
        "70 Backend-Tests + 12 Frontend-Tests aktuell.",
        "75 %-Coverage-Gate auf Backend (CI-blocking).",
        "Tests isolieren DB in tmp-File via conftest.py — kein Leak in die Demo-DB.",
        "LLM-Calls werden gemockt — keine Echtgeld-Calls in der Pipeline.",
        "Neue Features starten mit dem Test (TDD).",
        "ADMIN_API_KEY + Rate-Limit haben eigene Regression-Tests.",
    ], left=8.6, top=2.55, width=4.3, height=4.5, size=12)

    _add_footer(slide, n, total_pages)


def slide_cicd(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "CI / CD & Deployment", "Push to main → Pipeline → Render")

    stages = [
        ("ci.yml", "push / PR",
         "ruff + pytest (75 %) +\nVitest + ESLint +\nVite build + Playwright +\nbandit + pip-audit"),
        ("release.yml", "git tag v*.*.*", "Docker build →\nGHCR push →\nGitHub Release"),
        ("cd.yml", "CI grün auf main", "Render Deploy-Hook\n(API + Frontend)"),
    ]
    box_w = 4.0
    for i, (name, trigger, steps) in enumerate(stages):
        x = 0.6 + i * (box_w + 0.2)
        _add_rect(slide, left=x, top=2.0, width=box_w, height=3.6, fill=BG_WHITE, line=ACCENT)
        _add_text(slide, name, left=x + 0.2, top=2.1, width=box_w - 0.4, height=0.4,
                  size=17, color=ACCENT, bold=True, font="Menlo")
        _add_text(slide, "Trigger: " + trigger, left=x + 0.2, top=2.55, width=box_w - 0.4, height=0.4,
                  size=11, color=INK_MUTED)
        _add_text(slide, steps, left=x + 0.2, top=3.05, width=box_w - 0.4, height=2.4,
                  size=13, color=INK)
        # arrow to next stage
        if i < 2:
            ax = x + box_w + 0.02
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          Inches(ax), Inches(3.5), Inches(0.16), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()
            arrow.shadow.inherit = False

    # Deployment box
    _add_rect(slide, left=0.6, top=5.9, width=12.2, height=1.1, fill=ACCENT_SOFT, line=ACCENT)
    _add_text(slide, "Deployment", left=0.85, top=5.95, width=4, height=0.4,
              size=14, color=ACCENT, bold=True)
    _add_text(slide,
              "Render: Web-Service (FastAPI) + Static-Site (React) + Postgres   ·   "
              "Multi-stage Dockerfile (non-root, /health-Healthcheck)   ·   "
              "ADMIN_API_KEY-Gate · 30 req/min LLM-Rate-Limit · Forward-Merge-Migrations",
              left=0.85, top=6.35, width=11.8, height=0.6,
              size=11, color=INK_SOFT)

    _add_footer(slide, n, total_pages)


def slide_hardening(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Sicherheits-Härtung", "Vom Demo-PoC zum produktionstauglichen Setup")

    rows = [
        ("Vorher", "Nachher", "Belegt durch"),
        ("Admin-Endpoints offen für jedermann",
         "ADMIN_API_KEY-Pflicht · Default = 503 ohne Key",
         "tests/api/test_admin_auth.py · 5 Cases"),
        ("Stack-Traces im 500-Response",
         "Globaler Error-Handler · request_id-Korrelation",
         "tests/api/test_error_handlers.py"),
        ("LLM-Endpoints unbeschränkt rufbar",
         "30 req/min pro IP via slowapi",
         "tests/api/test_rate_limit.py"),
        ("CORS akzeptierte jede Render-App",
         "Regex eng auf eigene Frontend-Domain",
         "render.yaml"),
        ("Analytics-Recompute auf jedem Page-Load",
         "60 s TTL-Cache mit Invalidate auf Mutationen",
         "tests/unit/test_analytics_cache.py"),
        ("RAG-Index aktualisiert nur manuell",
         "Auto-Rebuild nach /tickets/decision",
         "SaveTriageDecisionUseCase"),
    ]
    box_x, box_y, box_w = 0.6, 2.0, 12.2
    col_w = [4.0, 4.5, 3.7]
    rh = 0.65
    for i, row in enumerate(rows):
        ry = box_y + i * rh
        bg = ACCENT_SOFT if i == 0 else (BG_WHITE if i % 2 == 1 else BG_PAGE)
        _add_rect(slide, left=box_x, top=ry, width=box_w, height=rh, fill=bg, line=BORDER)
        cx = box_x
        for j, cell in enumerate(row):
            _add_text(slide, cell, left=cx + 0.15, top=ry + 0.08, width=col_w[j] - 0.3,
                      height=rh - 0.1,
                      size=11, color=INK if i == 0 else INK_SOFT, bold=(i == 0))
            cx += col_w[j]

    _add_footer(slide, n, total_pages)


def slide_challenges(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Herausforderungen", "Was hat überraschend gebissen")

    items = [
        ("Postgres vs. SQLite FK-Verhalten",
         "Delete auf Tickets ohne vorheriges Delete der Events lief lokal grün, "
         "scheiterte auf Render mit ForeignKeyViolation. Fix: _delete_tickets()-Helper."),
        ("Test-Suite verschmutzte Demo-DB",
         "200+ pytest-Fixtures (WB-VIEWS-CLAUDIO, Workflow Close Test) landeten in der "
         "Live-DB. Fix: conftest.py setzt DATABASE_URL vor erstem app-Import auf tmp-File."),
        ("Composite-Prio kollabierte auf Score 9",
         "45 von 99 Tickets im Default-Bucket. Ursache: Matcher AND-verknüpfte tags_any "
         "und title_any. Fix: OR-Semantik + Default 2×2=4."),
        ("LocalStorage-Migration übersehen",
         "Returning-Users sahen die neuen Workbench-Spalten nicht. Fix: forward-merge "
         "fehlender Defaults in TicketsPage.jsx."),
        ("DB-Engine wird beim Import angelegt",
         "Macht Test-Isolation fragil. Aktuelle Lösung: Env-Var-Override vor Import. "
         "Sauberer wäre lazy get_engine() — bewusst nicht umgesetzt (Aufwand vs. Nutzen)."),
    ]
    y = 2.0
    for title, body in items:
        _add_rect(slide, left=0.6, top=y, width=12.2, height=0.95, fill=BG_WHITE, line=BORDER)
        _add_text(slide, title, left=0.85, top=y + 0.08, width=11.5, height=0.35,
                  size=14, color=ACCENT, bold=True)
        _add_text(slide, body, left=0.85, top=y + 0.42, width=11.7, height=0.6,
                  size=11, color=INK_SOFT)
        y += 1.02

    _add_footer(slide, n, total_pages)


def slide_lessons(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Lernerfahrungen", "Was ich beim nächsten Projekt mitnehme")

    items = [
        ("Retrieval schlägt Prompt-Engineering.",
         "Eine simple TF-IDF-Schicht über reviewte Tickets hat mehr Konsistenz gebracht "
         "als jede Prompt-Iteration."),
        ("Knowledge-Engineering und ML sind Komplemente.",
         "Der Klassifikator weiß was es ist; das Regelwerk weiß was zu tun ist. "
         "Beides zusammen ergibt explainable + auditable Triage."),
        ("Hexagonal zahlt sich exakt dann aus, wenn es weh tun würde.",
         "SQLite → Postgres war 1 Env-Var. LiteLLM → anderer Provider wäre 1 Datei. "
         "Disziplin am Anfang spart Stunden später."),
        ("AGENTS.md ist das wichtigste Artefakt im AI-Workflow.",
         "Jede User-Korrektur die als Regel kodifiziert wird, vermeidet identische Fehler "
         "in zukünftigen Sessions. Der Agent wird besser im Verlauf eines Projekts."),
        ("CI-Qualitätsgates sind die billigste Versicherung.",
         "Coverage-Gate, ruff, bandit, pip-audit, npm audit haben mehrfach Regressionen "
         "vor dem Merge gefangen."),
    ]
    y = 2.0
    for title, body in items:
        _add_rect(slide, left=0.6, top=y, width=12.2, height=0.95, fill=ACCENT_SOFT, line=ACCENT)
        _add_text(slide, title, left=0.85, top=y + 0.08, width=11.5, height=0.35,
                  size=15, color=ACCENT, bold=True)
        _add_text(slide, body, left=0.85, top=y + 0.42, width=11.7, height=0.6,
                  size=12, color=INK)
        y += 1.02

    _add_footer(slide, n, total_pages)


def slide_demo(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)
    _header(slide, "Live Demo", "AI-Empfehlung · KE-Priorisierung · Workbench · Reports")

    # Screenshot if available
    workbench = SCREENSHOT_DIR / "Tickets-workbench.png"
    if workbench.exists():
        slide.shapes.add_picture(str(workbench),
                                 Inches(0.6), Inches(2.0),
                                 width=Inches(8.0))
    _add_text(slide, "Probier es selbst", left=9.0, top=2.0, width=4, height=0.4,
              size=16, color=INK, bold=True)
    _add_bullets(slide, [
        "https://ai-assisted-ticket-triage-frontend.onrender.com/",
        "Sortier nach „Prio-Score“ → AML/Payments oben.",
        "Erfasse ein VPN-/Passwort-Ticket → siehe KE-Block im Modal.",
        "Reports → „Priorisierung & Aufwand“-Section.",
    ], left=9.0, top=2.55, width=4.0, height=3.5, size=11)

    _add_text(slide, "Repository", left=9.0, top=5.5, width=4, height=0.4,
              size=14, color=INK, bold=True)
    _add_text(slide, "github.com/clavinci94/ai-assisted-ticket-triage-platform",
              left=9.0, top=5.95, width=4, height=0.4,
              size=11, color=ACCENT)

    _add_footer(slide, n, total_pages)


def slide_qa(prs, total_pages, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG_PAGE)

    _add_text(slide, "Fragen?", left=1.0, top=3.0, width=11, height=1.5,
              size=72, color=ACCENT, bold=True)
    _add_text(slide,
              "Repo  ·  Live Demo  ·  AGENTS.md  ·  ADRs",
              left=1.0, top=4.4, width=11, height=0.6,
              size=22, color=INK_SOFT)
    _add_text(slide, "github.com/clavinci94/ai-assisted-ticket-triage-platform",
              left=1.0, top=5.1, width=11, height=0.5,
              size=14, color=INK_MUTED)

    _add_footer(slide, n, total_pages)


# ---------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------

def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    builders = [
        slide_title,
        slide_use_case,
        slide_architecture,
        slide_domain,
        slide_ai_layers,
        slide_ke_detail,
        slide_coding_agents,
        slide_agents_md,
        slide_tools,
        slide_testing,
        slide_cicd,
        slide_hardening,
        slide_challenges,
        slide_lessons,
        slide_demo,
        slide_qa,
    ]
    total = len(builders)

    for i, builder in enumerate(builders, start=1):
        if i == 1:
            builder(prs, total)
        else:
            builder(prs, total, i)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}")
